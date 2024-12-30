mport os
import time
import pandas as pd
import pytesseract
from PIL import Image
from googleapiclient.discovery import build
from google.auth.transport.requests import Request
from google.auth import default
from pptx import Presentation
import sqlite3
import pdfplumber
from flask import Flask, request, jsonify
from transformers import LlamaTokenizer, LlamaForCausalLM
import logging

# Setup Flask app
app = Flask(__name__)

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# ------------------ 1. Data Processing Functions ------------------

# 1.1 Process Excel File
def process_excel(file_path):
    start_time = time.time()
    try:
        excel_data = pd.read_excel(file_path, sheet_name=None)
        logger.info(f"Processed Excel in {time.time() - start_time:.4f} seconds")
        return excel_data
    except Exception as e:
        logger.error(f"Error processing Excel: {str(e)}")
        raise

# 1.2 Process Google Docs
def authenticate_google_api():
    creds, project = default()
    if not creds.valid:
        creds.refresh(Request())
    service = build('docs', 'v1', credentials=creds)
    return service

def process_google_doc(doc_id):
    start_time = time.time()
    try:
        service = authenticate_google_api()
        doc = service.documents().get(documentId=doc_id).execute()
        content = doc.get('body').get('content')
        
        text = ''
        for element in content:
            if 'paragraph' in element:
                for para in element['paragraph']['elements']:
                    text += para['textRun']['content']
        
        logger.info(f"Processed Google Doc in {time.time() - start_time:.4f} seconds")
        return text
    except Exception as e:
        logger.error(f"Error processing Google Doc: {str(e)}")
        raise

# 1.3 Process PDF
def process_pdf(file_path):
    start_time = time.time()
    try:
        with pdfplumber.open(file_path) as pdf:
            text = ''
            for page in pdf.pages:
                text += page.extract_text()
        logger.info(f"Processed PDF in {time.time() - start_time:.4f} seconds")
        return text
    except Exception as e:
        logger.error(f"Error processing PDF: {str(e)}")
        raise

# 1.4 Process Image/OCR
def process_image(image_path):
    start_time = time.time()
    try:
        img = Image.open(image_path)
        text = pytesseract.image_to_string(img)
        logger.info(f"Processed Image/OCR in {time.time() - start_time:.4f} seconds")
        return text
    except Exception as e:
        logger.error(f"Error processing Image/OCR: {str(e)}")
        raise

# 1.5 Process CDR (CSV format)
def process_cdr(file_path):
    start_time = time.time()
    try:
        cdr_df = pd.read_csv(file_path)
        cdr_df['call_time'] = pd.to_datetime(cdr_df['call_time'])
        cdr_df['time_of_day'] = cdr_df['call_time'].dt.hour.apply(lambda x: 'Day' if 6 <= x < 18 else 'Night')
        
        day_calls = cdr_df[cdr_df['time_of_day'] == 'Day']
        night_calls = cdr_df[cdr_df['time_of_day'] == 'Night']
        
        logger.info(f"Processed CDR in {time.time() - start_time:.4f} seconds")
        return day_calls, night_calls
    except Exception as e:
        logger.error(f"Error processing CDR: {str(e)}")
        raise

# 1.6 Process PowerPoint
def process_ppt(file_path):
    start_time = time.time()
    try:
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text
        logger.info(f"Processed PowerPoint in {time.time() - start_time:.4f} seconds")
        return text
    except Exception as e:
        logger.error(f"Error processing PowerPoint: {str(e)}")
        raise

# 1.7 Process SQL Database (SQLite)
def process_sql_db(db_path):
    start_time = time.time()
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        
        cursor.execute("SELECT * FROM calls")
        rows = cursor.fetchall()
        
        df = pd.DataFrame(rows, columns=[col[0] for col in cursor.description])
        conn.close()
        
        logger.info(f"Processed SQL Database in {time.time() - start_time:.4f} seconds")
        return df
    except Exception as e:
        logger.error(f"Error processing SQL DB: {str(e)}")
        raise


# ------------------ 2. Query LLaMA Model ------------------

def query_llama(data, model, tokenizer, inference_type="basic"):
    if inference_type == "basic":
        input_text = f"Analyze the following data: {data.head() if isinstance(data, pd.DataFrame) else data[:500]}"
    else:
        input_text = f"Perform advanced analysis: {data[:500]}"  # Adjust according to your requirements
    
    inputs = tokenizer(input_text, return_tensors="pt")
    
    start_time = time.time()
    outputs = model.generate(**inputs, max_length=512, temperature=0.7)
    result = tokenizer.decode(outputs[0], skip_special_tokens=True)
    
    logger.info(f"Query processed in {time.time() - start_time:.4f} seconds")
    return result


# ------------------ 3. Flask Route to Handle File Upload ------------------

@app.route('/process', methods=['POST'])
def process_file():
    file = request.files.get('file')
    inference_type = request.form.get('inference_type', 'basic')
    
    if not file:
        return jsonify({'error': 'No file provided'}), 400
    
    file_extension = file.filename.split('.')[-1].lower()
    
    # Save the file temporarily for processing
    temp_file_path = os.path.join('uploads', file.filename)
    file.save(temp_file_path)
    
    try:
        if file_extension == 'xlsx':
            data = process_excel(temp_file_path)
        elif file_extension == 'pdf':
            data = process_pdf(temp_file_path)
        elif file_extension in ['png', 'jpg']:
            data = process_image(temp_file_path)
        elif file_extension == 'csv':
            data = process_cdr(temp_file_path)
        elif file_extension == 'pptx':
            data = process_ppt(temp_file_path)
        elif file_extension == 'sql':
            data = process_sql_db(temp_file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400

        # Query LLaMA model with processed data
        insights = query_llama(data, model, tokenizer, inference_type)
        
        return jsonify({'insights': insights[:500]}), 200  # Return first 500 characters of insights
    
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ------------------ 4. Run the Flask App ------------------

if __name__ == '__main__':
    # Create directory for temporary uploads
    if not os.path.exists('uploads'):
        os.makedirs('uploads')
    
    model_name = "meta-llama/Llama-2-7b"
    tokenizer = LlamaTokenizer.from_pretrained(model_name)
    model = LlamaForCausalLM.from_pretrained(model_name)

    app.run(debug=False, host="0.0.0.0", port=5000)

