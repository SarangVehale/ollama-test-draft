mport time
import pandas as pd
import pytesseract
from PIL import Image
from googleapiclient.discovery import build
from google.auth.transport.requests import Request
from google.auth import default
from pptx import Presentation
import sqlite3
import pdfplumber
from transformers import LlamaTokenizer, LlamaForCausalLM


# ------------------ 1. Data Processing Functions ------------------

# 1.1 Process Excel File
def process_excel(file_path):
    start_time = time.time()
    excel_data = pd.read_excel(file_path, sheet_name=None)
    processing_time = time.time() - start_time
    print(f"Processed Excel in {processing_time:.4f} seconds")
    return excel_data

# 1.2 Process Google Docs
def authenticate_google_api():
    creds, project = default()
    if not creds.valid:
        creds.refresh(Request())
    service = build('docs', 'v1', credentials=creds)
    return service

def process_google_doc(doc_id):
    start_time = time.time()
    service = authenticate_google_api()
    doc = service.documents().get(documentId=doc_id).execute()
    content = doc.get('body').get('content')
    
    text = ''
    for element in content:
        if 'paragraph' in element:
            for para in element['paragraph']['elements']:
                text += para['textRun']['content']
    
    processing_time = time.time() - start_time
    print(f"Processed Google Doc in {processing_time:.4f} seconds")
    return text

# 1.3 Process PDF
def process_pdf(file_path):
    start_time = time.time()
    with pdfplumber.open(file_path) as pdf:
        text = ''
        for page in pdf.pages:
            text += page.extract_text()
    processing_time = time.time() - start_time
    print(f"Processed PDF in {processing_time:.4f} seconds")
    return text

# 1.4 Process Image/OCR
def process_image(image_path):
    start_time = time.time()
    img = Image.open(image_path)
    text = pytesseract.image_to_string(img)
    processing_time = time.time() - start_time
    print(f"Processed Image/OCR in {processing_time:.4f} seconds")
    return text

# 1.5 Process CDR (CSV format)
def process_cdr(file_path):
    start_time = time.time()
    cdr_df = pd.read_csv(file_path)
    cdr_df['call_time'] = pd.to_datetime(cdr_df['call_time'])
    cdr_df['time_of_day'] = cdr_df['call_time'].dt.hour.apply(lambda x: 'Day' if 6 <= x < 18 else 'Night')
    
    day_calls = cdr_df[cdr_df['time_of_day'] == 'Day']
    night_calls = cdr_df[cdr_df['time_of_day'] == 'Night']
    
    processing_time = time.time() - start_time
    print(f"Processed CDR in {processing_time:.4f} seconds")
    return day_calls, night_calls

# 1.6 Process PowerPoint
def process_ppt(file_path):
    start_time = time.time()
    prs = Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text
    processing_time = time.time() - start_time
    print(f"Processed PowerPoint in {processing_time:.4f} seconds")
    return text

# 1.7 Process SQL Database (SQLite)
def process_sql_db(db_path):
    start_time = time.time()
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    
    cursor.execute("SELECT * FROM calls")
    rows = cursor.fetchall()
    
    df = pd.DataFrame(rows, columns=[col[0] for col in cursor.description])
    conn.close()
    
    processing_time = time.time() - start_time
    print(f"Processed SQL Database in {processing_time:.4f} seconds")
    return df


# ------------------ 2. Query LLaMA Model ------------------

def query_llama(data, model, tokenizer, inference_type="basic"):
    if inference_type == "basic":
        input_text = f"Analyze the following data: {data.head() if isinstance(data, pd.DataFrame) else data[:500]}"
    else:
        input_text = f"Perform advanced analysis: {data[:500]}"  # Adjust according to your requirements
    
    inputs = tokenizer(input_text, return_tensors="pt")
    
    start_time = time.time()
    outputs = model.generate(**inputs, max_length=512, temperature=0.7)
    result = tokenizer.decode(outputs[0], skip_special_tokens=True)
    
    processing_time = time.time() - start_time
    print(f"Query processed in {processing_time:.4f} seconds")
    return result


# ------------------ 3. Orchestrator to Process and Query ------------------

def process_and_query(file_path, model, tokenizer, inference_type):
    file_extension = file_path.split('.')[-1].lower()
    start_time = time.time()

    if file_extension == 'xlsx':
        data = process_excel(file_path)
    elif file_extension == 'pdf':
        data = process_pdf(file_path)
    elif file_extension in ['png', 'jpg']:
        data = process_image(file_path)
    elif file_extension == 'csv':
        data = process_cdr(file_path)
    elif file_extension == 'pptx':
        data = process_ppt(file_path)
    elif file_extension == 'sql':
        data = process_sql_db(file_path)
    elif file_extension == 'gdoc':  # For Google Docs, use doc ID
        doc_id = file_path  # Assuming file_path is a doc ID
        data = process_google_doc(doc_id)
    else:
        print("Unsupported file type")
        return None

    processing_time = time.time() - start_time
    print(f"Total processing time for {file_extension} file: {processing_time:.4f} seconds")
    
    insights = query_llama(data, model, tokenizer, inference_type)
    return insights


# ------------------ 4. Main Execution with User Input ------------------

def get_user_input():
    print("Select the type of data you are processing:")
    print("1. Excel")
    print("2. Google Doc")
    print("3. PDF")
    print("4. Image/OCR")
    print("5. CDR (CSV)")
    print("6. PowerPoint")
    print("7. SQL Database")
    
    choice = int(input("Enter your choice (1-7): "))
    file_path = input("Enter the file path: ")

    print("\nSelect the type of inference you need:")
    print("1. Basic Analysis")
    print("2. Advanced Analysis")
    
    inference_choice = int(input("Enter your choice (1-2): "))
    inference_type = "basic" if inference_choice == 1 else "advanced"
    
    return choice, file_path, inference_type


def main():
    model_name = "meta-llama/Llama-2-7b"
    tokenizer = LlamaTokenizer.from_pretrained(model_name)
    model = LlamaForCausalLM.from_pretrained(model_name)

    # Get user input
    choice, file_path, inference_type = get_user_input()
    
    # Process the data and generate insights
    insights = process_and_query(file_path, model, tokenizer, inference_type)
    print(f"\nInsights Generated:\n{insights[:500]}...")  # Limit to first 500 characters


if __name__ == '__main__':
    main()

